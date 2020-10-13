import pandas as pd
import numpy as np
import PyPDF2
import docx2txt
import textract
from nltk.tokenize import word_tokenize
import re
import RAKE
import os
import operator
import string
from pptx import Presentation
import xlwt
from xlwt import Workbook
import xlrd

# Takes input from user
filename = input("Enter filename with extension (.txt, .doc, .ppt, .xls, .pdf): ")

# User can hard code PDF file, file path of raw text, and keywords
filepath = 'C:\\Users\\mguerra\\PyPrograms\\KeywordExtractor_RawText.txt'
keyword1 = "nerve"
keyword2 = "injury"
keyword3 = "repair"
keyword4 = "peripheral"
keyword5 = "damage"
keyword6 = "axon"
keyword7 = "neuroma"
keyword8 = "neuron"
keyword9 = "cell"
keyword10 = "tissue"

wb = Workbook()
sheet1 = wb.add_sheet('Results')
row = 1
count = 0
text = ""

#-----------------Converts Text Doc to Text------------------------

if(".txt" in filename):
    f = open(filename, "r")
    text = f.read()
    text = text.encode('ascii','ignore').lower()
    text = text.decode("utf-8")
    f.close()
    
#-----------------Converts Word Doc to Text------------------------

if(".doc" in filename or ".docx" in filename):
    text = docx2txt.process(filename)
    text = text.encode('ascii','ignore').lower()
    text = text.decode("utf-8")
    
#-----------------Converts PowerPoint to Text----------------------

if(".ppt" in filename or ".pptx" in filename):
    prs = Presentation(filename)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    text = "".join(text_runs)
    text = text.encode('ascii','ignore').lower()
    text = text.decode("utf-8")
    
#-----------------Converts Excel to Text---------------------------

if(".xls" in filename or ".xlsx" in filename):
    workbook = xlrd.open_workbook(filename)
    for i in range(workbook.nsheets):
        sheet = workbook.sheet_by_index(i)
        data = [[sheet.cell_value(r, c) for c in range(sheet.ncols)] for r in range(sheet.nrows)]
        text = "".join(str(v) for v in data)
        text = text.encode('ascii','ignore').lower()
        text = text.decode("utf-8")
        
#-----------------Converts PDF to Text-----------------------------

if(".pdf" in filename):
    pdfFileObj = open(filename,'rb')
    pdfReader = PyPDF2.PdfFileReader(pdfFileObj)
    num_pages = pdfReader.numPages

    while count < num_pages:
        pageObj = pdfReader.getPage(count)
        count +=1
        text += pageObj.extractText()

    if text != "":
        text = text
    else:
        text = textract.process('http://bit.ly/epo_keyword_extraction_document', method='tesseract', language='eng')

    text = text.encode('ascii','ignore').lower()
    text = text.decode("utf-8")
    tokens = word_tokenize(text)
    text = ""
    text = text.join(tokens)

#--------------------Converts Text to Excel------------------------

f = open(filepath, "w")
f.write(text)
f.close()

frequency = {}

with open(filepath, 'r') as myfile:
  text = myfile.read()

match_pattern = re.findall(r'\b[a-z]{3,15}\b', text)

for word in match_pattern:
    count = frequency.get(word,0)
    frequency[word] = count + 1
     
frequency_list = frequency.keys()
sorted_list = dict( sorted(frequency.items(), key=operator.itemgetter(1),reverse=True))
print("-----------------Keywords------------------")
sheet1.write(0, 0, filename)
#sheet1.write(1, 0, 'KEYWORDS')
#sheet1.write(1, 1, 'FREQUENCY')
# Search for Keyword 1
for words in frequency_list:
    if(words == keyword1):
        #row += 1
        #sheet1.write(row, 0, keyword1)
        #sheet1.write(row, 1, frequency[keyword1])
        print(keyword1, frequency[keyword1])

# Search for Keyword 2
for words in frequency_list:
    if(words == keyword2):
        #row += 1
        #sheet1.write(row, 0, keyword2)
        #sheet1.write(row, 1, frequency[keyword2])
        print(keyword2, frequency[keyword2])

# Search for Keyword 3
for words in frequency_list:
    if(words == keyword3):
        #row += 1
        #sheet1.write(row, 0, keyword3)
        #sheet1.write(row, 1, frequency[keyword3])
        print(keyword3, frequency[keyword3])

# Search for Keyword 4
for words in frequency_list:
    if(words == keyword4):
        #row += 1
        #sheet1.write(row, 0, keyword4)
        #sheet1.write(row, 1, frequency[keyword4])
        print(keyword4, frequency[keyword4])

# Search for Keyword 5
for words in frequency_list:
    if(words == keyword5):
        #row += 1
        #sheet1.write(row, 0, keyword5)
        #sheet1.write(row, 1, frequency[keyword5])
        print(keyword5, frequency[keyword5])

# Search for Keyword 6
for words in frequency_list:
    if(words == keyword6):
        #row += 1
        #sheet1.write(row, 0, keyword6)
        #sheet1.write(row, 1, frequency[keyword6])
        print(keyword6, frequency[keyword6])

# Search for Keyword 7
for words in frequency_list:
    if(words == keyword7):
        #row += 1
        #sheet1.write(row, 0, keyword7)
        #sheet1.write(row, 1, frequency[keyword7])
        print(keyword7, frequency[keyword7])

# Search for Keyword 8
for words in frequency_list:
    if(words == keyword8):
        #row += 1
        #sheet1.write(row, 0, keyword8)
        #sheet1.write(row, 1, frequency[keyword8])
        print(keyword8, frequency[keyword8])

# Search for Keyword 9
for words in frequency_list:
    if(words == keyword9):
        #row += 1
        #sheet1.write(row, 0, keyword9)
        #sheet1.write(row, 1, frequency[keyword9])
        print(keyword9, frequency[keyword9])

# Search for Keyword 10
for words in frequency_list:
    if(words == keyword10):
        #row += 1
        #sheet1.write(row, 0, keyword10)
        #sheet1.write(row, 1, frequency[keyword10])
        print(keyword10, frequency[keyword10])
        
print("---------------Other Words-----------------")
#row += 1
sheet1.write(row, 0, 'WORDS')
sheet1.write(row, 1, 'FREQUENCY')
for words in sorted_list:
    if(words == "the" or words == "and" or words == "for" or words == "from" or words == "that" or words == "with" or words == "are" or words == "not" or words == "than" or words == "only" or words == "after" or words == "into" or words == "was" or words == "any" or words == "this" or words == "can" or words == "also" or words == "more" or words == "these" or words == "who" or words == "had" or words == "were" or words == "has" or words == "its" or words == "which" or words == "your" or words == "whom" or words == "will" or words == "you" or words == "may" or words == "have" or words == "there" or words == "some" or words == "when" or words == "been" or words == "however" or words == "much" or words == "then" or words == "get" or words == "gets" or words == "where" or words == "about" or words == "like" or words == "they" or words == "within" or words == "per" or words == "other" or words == "make" or words == "come" or words == "comes" or words == "being" or words == "made" or words == "itself" or words == "while" or words == "even" or words == "without" or words == "going" or words == "own" or words == "them" or words == "what" or words == "just" or words == "yet" or words == "how" or words == "but" or words == "their" or words == "those"):
        pass
    else:
        row += 1
        sheet1.write(row, 0, words)
        sheet1.write(row, 1, frequency[words])
        print(words, frequency[words])

print("-------------------------------------------")
wb.save('KeywordExtractor_Results.xls')
